# pip install python-pptx
# python 하수처리장_관리대행_업체_현황_v2.py

"""
전국 하수처리장 관리대행 업체 현황 분석 보고서 자동 생성 스크립트
- 시각적 요소 강화 (차트, 인포그래픽, 카드 레이아웃)
- 전문적인 비즈니스 디자인 적용
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION

# --- 색상 정의 ---
COLOR_NAVY = RGBColor(0, 51, 102)       # #003366
COLOR_NAVY_DARK = RGBColor(0, 26, 51)   # #001a33
COLOR_GOLD = RGBColor(212, 175, 55)     # #d4af37
COLOR_GRAY_LIGHT = RGBColor(245, 245, 245) # #f5f5f5
COLOR_GRAY_BORDER = RGBColor(221, 221, 221) # #dddddd
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT_DARK = RGBColor(51, 51, 51)  # #333333
COLOR_BLUE_LIGHT = RGBColor(173, 216, 230)
COLOR_BLUE_MID = RGBColor(70, 130, 180)

def set_font(paragraph, font_name='Malgun Gothic', size=18, bold=False, color=COLOR_TEXT_DARK):
    """폰트 설정 유틸리티"""
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def create_title_slide(prs):
    """1장: 표지 (그라데이션 배경 효과, 골드 라인)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # 배경 사각형 (그라데이션 효과 대체: 상단 다크 네이비 -> 하단 네이비)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_NAVY
    bg.line.fill.background()
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(3.5), Cm(24), Cm(3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "전국 하수처리장 관리대행\n업체 현황 분석"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=36, bold=True, color=COLOR_WHITE)
    
    # 부제
    sub_box = slide.shapes.add_textbox(Cm(1), Cm(5.5), Cm(24), Cm(1.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2024 년 시장 동향 및 전망"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=20, color=RGBColor(204, 204, 204))
    
    # 날짜
    date_box = slide.shapes.add_textbox(Cm(1), Cm(6.5), Cm(24), Cm(1))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 년 07 월"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=16, color=RGBColor(204, 204, 204))
    
    # 골드 라인 장식
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(8), Cm(24), Cm(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_GOLD
    line.line.fill.background()

def create_toc_slide(prs):
    """2장: 목차 (2 단 그리드, 아이콘)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "목 차"
    set_font(p, size=28, bold=True, color=COLOR_NAVY)
    
    items = [
        "1. 시장 개요", "2. 주요 업체 현황", "3. 매출액 비교", "4. 지역별 분포", "5. 처리장 규모별 현황",
        "6. 입찰 동향", "7. 기술 개발 현황", "8. 시장 전망", "9. 결론"
    ]
    
    # 왼쪽 컬럼 (5 개)
    left_x = Cm(2)
    top_y = Cm(3)
    for i, text in enumerate(items[:5]):
        # 원형 번호 배경
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_x, top_y + i*Cm(1.2), Cm(0.8), Cm(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_NAVY
        circle.line.fill.background()
        # 번호 텍스트
        circ_tf = circle.text_frame
        circ_p = circ_tf.paragraphs[0]
        circ_p.text = str(i+1)
        circ_p.alignment = PP_ALIGN.CENTER
        set_font(circ_p, size=14, bold=True, color=COLOR_WHITE)
        
        # 항목 텍스트
        txt_box = slide.shapes.add_textbox(left_x + Cm(1.2), top_y + i*Cm(1.2), Cm(8), Cm(0.8))
        txt_tf = txt_box.text_frame
        txt_p = txt_tf.paragraphs[0]
        txt_p.text = text.split('. ')[1] if '. ' in text else text
        set_font(txt_p, size=18, color=COLOR_TEXT_DARK)

    # 오른쪽 컬럼 (4 개)
    right_x = Cm(10)
    for i, text in enumerate(items[5:]):
        idx = i + 6
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x, top_y + i*Cm(1.2), Cm(0.8), Cm(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_NAVY
        circle.line.fill.background()
        circ_tf = circle.text_frame
        circ_p = circ_tf.paragraphs[0]
        circ_p.text = str(idx)
        circ_p.alignment = PP_ALIGN.CENTER
        set_font(circ_p, size=14, bold=True, color=COLOR_WHITE)
        
        txt_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x + Cm(1.2), top_y + i*Cm(1.2), Cm(8), Cm(0.8))
        txt_tf = txt_box.text_frame
        txt_p = txt_tf.paragraphs[0]
        txt_p.text = text.split('. ')[1] if '. ' in text else text
        set_font(txt_p, size=18, color=COLOR_TEXT_DARK)

def create_market_overview_slide(prs):
    """3 장: 시장 개요 (인포그래픽 카드)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "시장 개요"
    set_font(p, size=28, bold=True, color=COLOR_NAVY)
    
    # 3 개의 카드 데이터
    cards = [
        {"value": "2,458 개소", "label": "전국 하수처리장 총량", "desc": "2024 년 기준"},
        {"value": "67.3%", "label": "민간위탁 비율", "desc": "지속적 증가 추세"},
        {"value": "4.2%", "label": "연평균 성장률 (CAGR)", "desc": "2020-2024 년"}
    ]
    
    card_width = Cm(8.5)
    card_height = Cm(5)
    gap = Cm(0.5)
    start_x = Cm(1.5)
    start_y = Cm(2.5)
    
    for i, card in enumerate(cards):
        x = start_x + i * (card_width + gap)
        
        # 카드 배경
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_width, card_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLOR_GRAY_LIGHT
        rect.line.color.rgb = COLOR_NAVY
        rect.line.width = Pt(2)
        
        # 값 (큰 숫자)
        val_box = slide.shapes.add_textbox(x + Cm(0.5), start_y + Cm(0.5), card_width - Cm(1), Cm(1.5))
        val_tf = val_box.text_frame
        val_p = val_tf.paragraphs[0]
        val_p.text = card["value"]
        val_p.alignment = PP_ALIGN.CENTER
        set_font(val_p, size=48, bold=True, color=COLOR_NAVY)
        
        # 라벨
        label_box = slide.shapes.add_textbox(x + Cm(0.5), start_y + Cm(2.2), card_width - Cm(1), Cm(0.8))
        label_tf = label_box.text_frame
        label_p = label_tf.paragraphs[0]
        label_p.text = card["label"]
        label_p.alignment = PP_ALIGN.CENTER
        set_font(label_p, size=18, bold=True, color=COLOR_TEXT_DARK)
        
        # 설명
        desc_box = slide.shapes.add_textbox(x + Cm(0.5), start_y + Cm(3.2), card_width - Cm(1), Cm(0.8))
        desc_tf = desc_box.text_frame
        desc_p = desc_tf.paragraphs[0]
        desc_p.text = card["desc"]
        desc_p.alignment = PP_ALIGN.CENTER
        set_font(desc_p, size=14, color=RGBColor(100, 100, 100))
    
    # 하단 시장 규모 텍스트
    market_box = slide.shapes.add_textbox(Cm(1), Cm(8.5), Cm(26), Cm(1))
    m_tf = market_box.text_frame
    m_p = m_tf.paragraphs[0]
    m_p.text = "시장 규모: 약 2 조 8,500 억원 (2024 년 예상)"
    m_p.alignment = PP_ALIGN.CENTER
    set_font(m_p, size=20, bold=True, color=COLOR_NAVY)

def create_main_companies_slide(prs):
    """4 장: 주요 업체 현황 (개선된 표)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "주요 업체 현황 (상위 8 개사)"
    set_font(p, size=28, bold=True, color=COLOR_NAVY)
    
    # 데이터
    data = [
        ["업체명", "운영 처리장 수", "처리용량 (톤/일)", "시장점유율 (%)"],
        ["A 환경그룹", "156", "2,850,000", "18.5"],
        ["B 워터테크", "128", "2,100,000", "14.2"],
        ["C 그린시스템", "98", "1,780,000", "11.8"],
        ["D 에코서비스", "87", "1,650,000", "10.5"],
        ["E 클린워터", "72", "1,320,000", "8.9"],
        ["F 물관리공사", "65", "1,180,000", "7.7"],
        ["G 환경산업", "58", "980,000", "6.5"],
        ["H 워터솔루션", "45", "820,000", "5.4"],
    ]
    
    # 표 생성
    rows = len(data)
    cols = 4
    left = Cm(1.5)
    top = Cm(2.5)
    width = Cm(25)
    height = Cm(6)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 열 너비 조정
    table.columns[0].width = Cm(6)
    table.columns[1].width = Cm(5)
    table.columns[2].width = Cm(7)
    table.columns[3].width = Cm(7)
    
    # 데이터 채우기 및 스타일링
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 폰트 설정
            for paragraph in cell.text_frame.paragraphs:
                if r == 0: # 헤더
                    set_font(paragraph, size=14, bold=True, color=COLOR_WHITE)
                else:
                    set_font(paragraph, size=14, color=COLOR_TEXT_DARK)
            
            # 셀 서식
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_NAVY
            else:
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_WHITE
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_GRAY_LIGHT
            
            # 테두리 (python-pptx 는 셀 개별 테두리 설정을 지원하지 않음 - 표 전체 테두리로 대체)
            # 표 생성 시 자동으로 회색 테두리가 적용됨
            pass

    # 하단 강조 텍스트
    summary_box = slide.shapes.add_textbox(Cm(1.5), Cm(8.8), Cm(25), Cm(1))
    s_tf = summary_box.text_frame
    s_p = s_tf.paragraphs[0]
    s_p.text = "상위 8 개사 합계 시장점유율: 83.5%"
    s_p.alignment = PP_ALIGN.RIGHT
    set_font(s_p, size=16, bold=True, color=COLOR_GOLD)

def create_revenue_chart_slide(prs):
    """5 장: 매출액 현황 (수평 막대 차트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "매출액 현황 (최근 3 개년 상위 5 개사)"
    set_font(p, size=28, bold=True, color=COLOR_NAVY)
    
    # 데이터 (단위: 억원)
    companies = ["A 환경그룹", "B 워터테크", "C 그린시스템", "D 에코서비스", "E 클린워터"]
    data_2022 = [3200, 2850, 2100, 1950, 1600]
    data_2023 = [3450, 3050, 2280, 2100, 1750]
    data_2024 = [3680, 3250, 2450, 2250, 1890]
    
    chart_left = Cm(2)
    chart_top = Cm(2.5)
    chart_width = Cm(24)
    chart_height = Cm(6)
    
    # 차트 데이터 준비
    chart_data = CategoryChartData()
    chart_data.categories = companies
    chart_data.add_series('2022 년', data_2022)
    chart_data.add_series('2023 년', data_2023)
    chart_data.add_series('2024 년', data_2024)
    
    # 차트 추가 (Clustered Bar Chart)
    x, y, cx, cy = chart_left, chart_top, chart_width, chart_height
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    
    # 차트 스타일링
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = True
    # 그리드라인 색상 설정 (python-pptx 제한으로 생략)
    
    # 시리즈 색상 설정
    series_colors = [COLOR_BLUE_LIGHT, COLOR_BLUE_MID, COLOR_NAVY]
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = series_colors[i]
        series.format.fill.transparency = 0.1
        # 데이터 라벨 추가
        series.has_data_labels = True
        series.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        series.data_labels.number_format = '#,##0'
        series.data_labels.font.size = Pt(10)
        series.data_labels.font.color.rgb = COLOR_TEXT_DARK

def create_region_map_slide(prs):
    """6 장: 지역별 업체 분포 (지도 개요 + 표)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "지역별 업체 분포"
    set_font(p, size=28, bold=True, color=COLOR_NAVY)
    
    # 왼쪽: 간단한 지도 개요 (도형 조합)
    map_left = Cm(1.5)
    map_top = Cm(2.5)
    
    # 강원권
    shape_gw = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, map_left, map_top, Cm(4), Cm(2.5))
    shape_gw.fill.solid()
    shape_gw.fill.fore_color.rgb = COLOR_BLUE_LIGHT
    shape_gw.line.color.rgb = COLOR_WHITE
    gw_tf = shape_gw.text_frame
    gw_tf.paragraphs[0].text = "강원권\n89 개"
    set_font(gw_tf.paragraphs[0], size=12, bold=True, color=COLOR_WHITE)
    
    # 수도권 (경기/서울/인천) - 사각형
    shape_seoul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, map_left + Cm(4.2), map_top + Cm(0.5), Cm(3.5), Cm(2))
    shape_seoul.fill.solid()
    shape_seoul.fill.fore_color.rgb = COLOR_NAVY
    shape_seoul.line.color.rgb = COLOR_WHITE
    seoul_tf = shape_seoul.text_frame
    seoul_tf.paragraphs[0].text = "수도권\n485 개"
    set_font(seoul_tf.paragraphs[0], size=12, bold=True, color=COLOR_WHITE)
    
    # 충청권
    shape_chung = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, map_left + Cm(2), map_top + Cm(2.8), Cm(4), Cm(2))
    shape_chung.fill.solid()
    shape_chung.fill.fore_color.rgb = COLOR_BLUE_MID
    shape_chung.line.color.rgb = COLOR_WHITE
    chung_tf = shape_chung.text_frame
    chung_tf.paragraphs[0].text = "충청권\n198 개"
    set_font(chung_tf.paragraphs[0], size=12, bold=True, color=COLOR_WHITE)
    
    # 호남권
    shape_honam = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, map_left + Cm(1.5), map_top + Cm(5), Cm(3.5), Cm(2.5))
    shape_honam.fill.solid()
    shape_honam.fill.fore_color.rgb = COLOR_BLUE_LIGHT
    shape_honam.line.color.rgb = COLOR_WHITE
    honam_tf = shape_honam.text_frame
    honam_tf.paragraphs[0].text = "호남권\n245 개"
    set_font(honam_tf.paragraphs[0], size=12, bold=True, color=COLOR_WHITE)
    
    # 영남권
    shape_yeong = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, map_left + Cm(5.5), map_top + Cm(3), Cm(4.5), Cm(4))
    shape_yeong.fill.solid()
    shape_yeong.fill.fore_color.rgb = COLOR_NAVY_DARK
    shape_yeong.line.color.rgb = COLOR_WHITE
    yeong_tf = shape_yeong.text_frame
    yeong_tf.paragraphs[0].text = "영남권\n312 개"
    set_font(yeong_tf.paragraphs[0], size=12, bold=True, color=COLOR_WHITE)
    
    # 기타 (제주 등)
    shape_etc = slide.shapes.add_shape(MSO_SHAPE.OVAL, map_left + Cm(1), map_top + Cm(8), Cm(2), Cm(1))
    shape_etc.fill.solid()
    shape_etc.fill.fore_color.rgb = COLOR_GRAY_BORDER
    shape_etc.line.color.rgb = COLOR_WHITE
    etc_tf = shape_etc.text_frame
    etc_tf.paragraphs[0].text = "기타\n169 개"
    set_font(etc_tf.paragraphs[0], size=10, bold=True, color=COLOR_TEXT_DARK)
    
    # 오른쪽: 표
    table_data = [
        ["권역", "업체 수", "비중 (%)"],
        ["수도권 (경기/서울/인천)", "485", "32.5"],
        ["영남권 (부산/대구/경남/울산)", "312", "20.9"],
        ["호남권 (광주/전남/전북/제주)", "245", "16.4"],
        ["충청권 (대전/세종/충남/충북)", "198", "13.3"],
        ["강원권", "89", "5.9"],
        ["기타", "169", "11.0"]
    ]
    
    rows = len(table_data)
    cols = 3
    t_left = Cm(11)
    t_top = Cm(2.5)
    t_width = Cm(16)
    t_height = Cm(6)
    
    table = slide.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height).table
    table.columns[0].width = Cm(8)
    table.columns[1].width = Cm(4)
    table.columns[2].width = Cm(4)
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = table_data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                if r == 0:
                    set_font(paragraph, size=14, bold=True, color=COLOR_WHITE)
                else:
                    set_font(paragraph, size=14, color=COLOR_TEXT_DARK)
            
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE if r % 2 == 0 else COLOR_GRAY_LIGHT
            
            # 테두리 (python-pptx 는 셀 개별 테두리 설정을 지원하지 않음)
            pass

def create_scale_pie_slide(prs):
    """7 장: 처리장 규모별 현황 (도넛 차트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "처리장 규모별 현황"
    set_font(p, size=28, bold=True, color=COLOR_NAVY)
    
    # 데이터
    categories = ["대규모\n(10,000 톤/일 이상)", "중규모\n(1,000-10,000 톤/일)", "소규모\n(1,000 톤/일 미만)"]
    values = [15, 45, 40]
    
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('비중', values)
    
    # 차트 위치
    x, y, cx, cy = Cm(2), Cm(2.5), Cm(12), Cm(7)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    ).chart
    
    # 스타일링
    chart.has_legend = False # 범례는 수동으로 옆에 배치
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.show_percentage = True
    chart.plots[0].data_labels.show_value = False
    chart.plots[0].data_labels.font.size = Pt(14)
    chart.plots[0].data_labels.font.bold = True
    chart.plots[0].data_labels.font.color.rgb = COLOR_TEXT_DARK
    
    # 슬라이스 색상
    colors = [COLOR_NAVY, COLOR_BLUE_MID, COLOR_BLUE_LIGHT]
    for i, point in enumerate(chart.plots[0].series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[i]
    
    # 수동 범례 (오른쪽)
    legend_x = Cm(15)
    legend_y = Cm(3)
    labels = [
        ("대규모 (10k 톤+/일)", "15%", COLOR_NAVY),
        ("중규모 (1k-10k 톤/일)", "45%", COLOR_BLUE_MID),
        ("소규모 (1k 톤-/일)", "40%", COLOR_BLUE_LIGHT)
    ]
    
    for i, (text, pct, color) in enumerate(labels):
        # 색상 박스
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_x, legend_y + i*Cm(1.2), Cm(0.5), Cm(0.5))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # 텍스트
        txt = slide.shapes.add_textbox(legend_x + Cm(0.8), legend_y + i*Cm(1.2), Cm(10), Cm(0.8))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = f"{text} : {pct}"
        set_font(p, size=16, color=COLOR_TEXT_DARK)

def create_bid_trend_slide(prs):
    """8 장: 입찰 동향 (라인 차트)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "입찰 동향 (최근 5 개년)"
    set_font(p, size=28, bold=True, color=COLOR_NAVY)
    
    # 데이터
    years = ['2020', '2021', '2022', '2023', '2024']
    amounts = [420, 485, 520, 580, 625] # 억원
    
    chart_data = CategoryChartData()
    chart_data.categories = years
    chart_data.add_series('평균 입찰 금액 (억원)', amounts)
    
    x, y, cx, cy = Cm(2), Cm(2.5), Cm(24), Cm(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
    ).chart
    
    # 스타일링
    chart.has_legend = False
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = True
    # 그리드라인 색상 설정 (python-pptx 제한으로 생략)
    
    series = chart.series[0]
    series.format.line.width = Pt(3)
    # 라인 색상 설정 (python-pptx 제한으로 fill 사용)
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = COLOR_GOLD
    series.marker.style = 8 # Diamond
    series.marker.size = 10
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = COLOR_NAVY
    
    # 데이터 라벨
    series.has_data_labels = True
    series.data_labels.position = XL_DATA_LABEL_POSITION.ABOVE
    series.data_labels.font.size = Pt(14)
    series.data_labels.font.bold = True
    series.data_labels.font.color.rgb = COLOR_NAVY
    series.data_labels.number_format = '0" 억"'
    
    # 하단 낙찰률 강조
    rate_box = slide.shapes.add_textbox(Cm(2), Cm(8.5), Cm(24), Cm(1))
    rt_tf = rate_box.text_frame
    rt_p = rt_tf.paragraphs[0]
    rt_p.text = "평균 낙찰률: 92.3% (안정적인 시장 형성)"
    rt_p.alignment = PP_ALIGN.CENTER
    set_font(rt_p, size=20, bold=True, color=COLOR_NAVY)

def create_tech_trend_slide(prs):
    """9 장: 기술 개발 현황 (3 단 카드 레이아웃)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "기술 개발 현황"
    set_font(p, size=28, bold=True, color=COLOR_NAVY)
    
    techs = [
        {
            "title": "AI 기반 운영 최적화",
            "icon": "🤖", # 텍스트로 아이콘 대체 (Shape 복잡성回避)
            "desc1": "에너지 절감 25%",
            "desc2": "운영 효율 30% 향상"
        },
        {
            "title": "에너지 자립화",
            "icon": "⚡",
            "desc1": "태양광 + 바이오가스",
            "desc2": "자립률 45% 달성"
        },
        {
            "title": "슬러지 감량화",
            "icon": "♻️",
            "desc1": "부피 60% 감소",
            "desc2": "재활용률 80%"
        }
    ]
    
    card_w = Cm(8.5)
    card_h = Cm(5.5)
    gap = Cm(0.5)
    start_x = Cm(1.5)
    start_y = Cm(2.5)
    
    for i, tech in enumerate(techs):
        x = start_x + i * (card_w + gap)
        
        # 카드 배경
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_w, card_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLOR_WHITE
        rect.line.color.rgb = COLOR_NAVY
        rect.line.width = Pt(2)
        
        # 아이콘 (원형 배경)
        icon_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Cm(0.5), start_y + Cm(0.5), Cm(1.2), Cm(1.2))
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = COLOR_GOLD
        icon_circle.line.fill.background()
        icon_tf = icon_circle.text_frame
        icon_p = icon_tf.paragraphs[0]
        icon_p.text = tech["icon"]
        icon_p.alignment = PP_ALIGN.CENTER
        set_font(icon_p, size=20, color=COLOR_WHITE)
        
        # 제목
        title_box = slide.shapes.add_textbox(x + Cm(2), start_y + Cm(0.5), card_w - Cm(2.5), Cm(1))
        t_tf = title_box.text_frame
        t_p = t_tf.paragraphs[0]
        t_p.text = tech["title"]
        set_font(t_p, size=18, bold=True, color=COLOR_NAVY)
        
        # 설명 1
        d1_box = slide.shapes.add_textbox(x + Cm(0.5), start_y + Cm(2), card_w - Cm(1), Cm(0.8))
        d1_tf = d1_box.text_frame
        d1_p = d1_tf.paragraphs[0]
        d1_p.text = tech["desc1"]
        set_font(d1_p, size=16, color=COLOR_TEXT_DARK)
        
        # 설명 2
        d2_box = slide.shapes.add_textbox(x + Cm(0.5), start_y + Cm(3), card_w - Cm(1), Cm(0.8))
        d2_tf = d2_box.text_frame
        d2_p = d2_tf.paragraphs[0]
        d2_p.text = tech["desc2"]
        set_font(d2_p, size=16, color=COLOR_TEXT_DARK)

def create_conclusion_slide(prs):
    """10 장: 결론 및 제언"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(Cm(1), Cm(1.5), Cm(26), Cm(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "결론 및 제언"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=32, bold=True, color=COLOR_NAVY)
    
    # 결론 리스트
    conclusions = [
        "시장 안정성: 연 4.2% 성장으로 안정적 투자처",
        "기술 혁신: AI·에너지 자립화 기술이 경쟁력 결정",
        "투자 유망 분야: 스마트 워터관리, 에너지 효율화, 슬러지 자원화"
    ]
    
    box_top = Cm(4)
    for i, text in enumerate(conclusions):
        # 불릿 포인트 배경 (작은 원)
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(5), box_top + i*Cm(1.5), Cm(0.4), Cm(0.4))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = COLOR_GOLD
        bullet.line.fill.background()
        
        # 텍스트
        txt_box = slide.shapes.add_textbox(Cm(6), box_top + i*Cm(1.5), Cm(20), Cm(1))
        t_tf = txt_box.text_frame
        t_p = t_tf.paragraphs[0]
        t_p.text = text
        set_font(t_p, size=22, color=COLOR_TEXT_DARK)
    
    # 감사 인사
    thank_box = slide.shapes.add_textbox(Cm(1), Cm(9), Cm(26), Cm(1.5))
    th_tf = thank_box.text_frame
    th_p = th_tf.paragraphs[0]
    th_p.text = "감사합니다"
    th_p.alignment = PP_ALIGN.CENTER
    set_font(th_p, size=24, bold=True, color=COLOR_NAVY)
    
    # 하단 골드 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(10.2), Cm(24), Cm(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_GOLD
    line.line.fill.background()

def main():
    try:
        prs = Presentation()
        # 슬라이드 크기 설정 (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 슬라이드 생성
        create_title_slide(prs)
        create_toc_slide(prs)
        create_market_overview_slide(prs)
        create_main_companies_slide(prs)
        create_revenue_chart_slide(prs)
        create_region_map_slide(prs)
        create_scale_pie_slide(prs)
        create_bid_trend_slide(prs)
        create_tech_trend_slide(prs)
        create_conclusion_slide(prs)
        
        # 저장
        output_filename = '하수처리장_관리대행_업체_현황.pptx'
        prs.save(output_filename)
        print(f"✅ 성공적으로 파일이 생성되었습니다: {output_filename}")
        
    except Exception as e:
        print(f"❌ 오류 발생: {e}")

if __name__ == "__main__":
    main()
