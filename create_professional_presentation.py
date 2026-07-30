#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
전국 하수처리장 관리대행 업체 현황 분석 - 전문 비즈니스 프레젠테이션
화려하고 전문적인 디자인의 파워포인트 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import matplotlib.pyplot as plt
import io

# ============================================================================
# 1. 상수 정의 (색상, 폰트 크기 등)
# ============================================================================

# 색상 팔레트
COLOR_NAVY = RGBColor(27, 58, 92)        # #1B3A5C - 메인 컬러
COLOR_NAVY_DARK = RGBColor(15, 37, 64)   # #0F2540 - 그라데이션용
COLOR_GOLD = RGBColor(212, 175, 55)      # #D4AF37 - 강조용
COLOR_BLUE = RGBColor(74, 144, 226)      # #4A90E2 - 액센트/차트용
COLOR_WHITE = RGBColor(255, 255, 255)    # #FFFFFF - 배경
COLOR_LIGHT_GRAY = RGBColor(248, 249, 250)  # #F8F9FA - 배경
COLOR_TEXT = RGBColor(44, 62, 80)        # #2C3E50 - 텍스트
COLOR_TEXT_LIGHT = RGBColor(127, 140, 141)  # #7F8C8D - 부제목
COLOR_LIGHT_BLUE = RGBColor(235, 245, 251)  # #EBF5FB - 표 행/박스
COLOR_CHART_2022 = RGBColor(174, 214, 241)  # #AED6F1
COLOR_CHART_2023 = RGBColor(93, 173, 226)   # #5DADE2
COLOR_CHART_2024 = RGBColor(27, 58, 92)     # #1B3A5C

# 폰트 크기
FONT_SIZE_TITLE = Pt(32)
FONT_SIZE_SUBTITLE = Pt(20)
FONT_SIZE_BODY = Pt(16)
FONT_SIZE_TABLE_HEADER = Pt(14)
FONT_SIZE_TABLE_BODY = Pt(12)
FONT_SIZE_CARD_NUMBER = Pt(48)
FONT_SIZE_CARD_TEXT = Pt(14)

# 슬라이드 설정
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_LEFT = Inches(1.5)
MARGIN_RIGHT = Inches(1.5)
MARGIN_TOP = Inches(1)
MARGIN_BOTTOM = Inches(1)


# ============================================================================
# 2. 유틸리티 함수
# ============================================================================

def set_paragraph_format(paragraph, font_size=FONT_SIZE_BODY, bold=False, 
                         color=COLOR_TEXT, align=PP_ALIGN.LEFT):
    """문단 서식 설정"""
    paragraph.font.size = font_size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.alignment = align
    try:
        paragraph.font.name = '맑은 고딕'
    except:
        pass

def create_title_slide(prs, title, subtitle, date):
    """표지 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경: 네이비 블루 그라데이션 (단색으로 표현)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_NAVY
    bg.line.fill.background()
    
    # 제목 (중앙 상단)
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.25), SLIDE_WIDTH - MARGIN_LEFT * 2, Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_format(p, font_size=Pt(36), bold=True, color=COLOR_WHITE)
    
    # 그림자 효과 (간단한 방법으로 구현)
    title_box.shadow.inherit = False
    
    # 부제 (제목 바로 아래)
    subtitle_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(3.5), SLIDE_WIDTH - MARGIN_LEFT * 2, Inches(0.8)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_format(p, font_size=FONT_SIZE_SUBTITLE, bold=False, color=COLOR_GOLD)
    
    # 날짜 (중앙 하단)
    date_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(6), SLIDE_WIDTH - MARGIN_LEFT * 2, Inches(0.5)
    )
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = date
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_format(p, font_size=Pt(16), bold=False, 
                        color=RGBColor(180, 180, 180))
    
    # 장식: 하단에 골드 라인
    gold_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        SLIDE_WIDTH * 0.2, Inches(6.8), 
        SLIDE_WIDTH * 0.6, Pt(3)
    )
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = COLOR_GOLD
    gold_line.line.fill.background()
    
    return slide

def create_toc_slide(prs, items):
    """목차 슬라이드 생성 (3x3 그리드)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, Inches(4), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "목차"
    set_paragraph_format(p, font_size=Pt(28), bold=True, color=COLOR_NAVY)
    
    # 3x3 그리드 레이아웃
    grid_start_x = MARGIN_LEFT
    grid_start_y = Inches(2)
    grid_width = (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT) / 3
    grid_height = Inches(1.2)
    
    for idx, item in enumerate(items):
        row = idx // 3
        col = idx % 3
        
        x = grid_start_x + col * grid_width
        y = grid_start_y + row * grid_height
        
        # 번호 원형 배경
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y + Inches(0.1), Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_NAVY
        circle.line.fill.background()
        
        # 번호 텍스트
        num_box = slide.shapes.add_textbox(
            x, y + Inches(0.1), Inches(0.6), Inches(0.6)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(idx + 1)
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_paragraph_format(p, font_size=Pt(24), bold=True, color=COLOR_WHITE)
        
        # 항목 텍스트
        text_box = slide.shapes.add_textbox(
            x + Inches(0.8), y + Inches(0.15), grid_width - Inches(1), Inches(0.6)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        set_paragraph_format(p, font_size=Pt(16), bold=False, color=COLOR_TEXT)
    
    return slide

def create_infographic_cards_slide(prs, title, cards, footer_text=None):
    """인포그래픽 카드 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, Inches(6), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_format(p, font_size=Pt(28), bold=True, color=COLOR_NAVY)
    
    # 카드 레이아웃
    card_width = (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.4)) / 3
    card_height = Inches(3.5)
    card_start_y = Inches(2)
    
    for idx, card in enumerate(cards):
        x = MARGIN_LEFT + idx * (card_width + Inches(0.2))
        
        # 카드 배경
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, card_start_y, card_width, card_height
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = COLOR_WHITE
        card_shape.line.color.rgb = RGBColor(200, 200, 200)
        card_shape.line.width = Pt(1)
        
        # 큰 숫자
        num_box = slide.shapes.add_textbox(
            x + Inches(0.3), card_start_y + Inches(0.5), 
            card_width - Inches(0.6), Inches(1.5)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = card['number']
        set_paragraph_format(p, font_size=FONT_SIZE_CARD_NUMBER, 
                           bold=True, color=COLOR_NAVY)
        p.alignment = PP_ALIGN.CENTER
        
        # 설명 텍스트
        desc_box = slide.shapes.add_textbox(
            x + Inches(0.3), card_start_y + Inches(2), 
            card_width - Inches(0.6), Inches(1.2)
        )
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = card['description']
        set_paragraph_format(p, font_size=FONT_SIZE_CARD_TEXT, 
                           bold=False, color=COLOR_TEXT_LIGHT)
        p.alignment = PP_ALIGN.CENTER
    
    # 푸터 박스
    if footer_text:
        footer_y = card_start_y + card_height + Inches(0.3)
        footer = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            MARGIN_LEFT, footer_y, 
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.8)
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = COLOR_LIGHT_BLUE
        footer.line.fill.background()
        
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_format(p, font_size=Pt(16), bold=True, color=COLOR_NAVY)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def create_table_slide(prs, title, subtitle, headers, data, highlight_text=None):
    """표 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, Inches(6), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_format(p, font_size=Pt(28), bold=True, color=COLOR_NAVY)
    
    # 부제
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            MARGIN_LEFT, MARGIN_TOP + Inches(0.7), Inches(6), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        set_paragraph_format(p, font_size=Pt(16), bold=False, color=COLOR_TEXT_LIGHT)
    
    # 표 생성
    table_rows = len(data) + 1
    table_cols = len(headers)
    table_width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    table_height = Inches(0.4) * table_rows + Inches(0.1)
    table_y = Inches(2.2)
    
    table_shape = slide.shapes.add_table(
        table_rows, table_cols, MARGIN_LEFT, table_y, table_width, table_height
    )
    table = table_shape.table
    
    # 헤더 설정
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_format(p, font_size=FONT_SIZE_TABLE_HEADER, 
                           bold=True, color=COLOR_WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
    
    # 데이터 행 설정
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            
            # 숫자는 오른쪽 정렬
            if col_idx > 0:
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.CENTER
            
            set_paragraph_format(p, font_size=FONT_SIZE_TABLE_BODY, 
                               bold=False, color=COLOR_TEXT)
            
            # 번갈아 배경색
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    
    # 강조 박스
    if highlight_text:
        highlight_y = table_y + table_height + Inches(0.3)
        highlight = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            MARGIN_LEFT, highlight_y, 
            table_width, Inches(0.6)
        )
        highlight.fill.solid()
        highlight.fill.fore_color.rgb = COLOR_GOLD
        highlight.line.fill.background()
        
        tf = highlight.text_frame
        p = tf.paragraphs[0]
        p.text = highlight_text
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_format(p, font_size=Pt(14), bold=True, color=COLOR_WHITE)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def create_bar_chart_slide(prs, title, subtitle, categories, series_data, series_names):
    """수평 막대 차트 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, Inches(6), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_format(p, font_size=Pt(28), bold=True, color=COLOR_NAVY)
    
    # 부제
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            MARGIN_LEFT, MARGIN_TOP + Inches(0.7), Inches(8), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        set_paragraph_format(p, font_size=Pt(16), bold=False, color=COLOR_TEXT_LIGHT)
    
    # matplotlib로 차트 생성
    fig, ax = plt.subplots(figsize=(10, 5))
    
    bar_width = 0.25
    y_positions = list(range(len(categories)))
    
    colors = ['#AED6F1', '#5DADE2', '#1B3A5C']
    
    for i, (name, data) in enumerate(zip(series_names, series_data)):
        offsets = [y + i * bar_width - bar_width for y in y_positions]
        bars = ax.barh(offsets, data, height=bar_width, color=colors[i], label=name)
        
        # 막대 위에 숫자 표시
        for j, val in enumerate(data):
            ax.text(val + 50, offsets[j], 
                   f'{val:,}', va='center', fontsize=10, color='#2C3E50')
    
    ax.set_yticks(y_positions)
    ax.set_yticklabels(categories, fontsize=12)
    ax.set_xlabel('매출액 (억원)', fontsize=12)
    ax.legend(loc='upper right', fontsize=10)
    ax.grid(axis='x', linestyle='--', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    
    plt.tight_layout()
    
    # 차트를 이미지로 저장
    img_buffer = io.BytesIO()
    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close()
    
    # 슬라이드에 추가
    slide.shapes.add_picture(
        img_buffer, MARGIN_LEFT, Inches(1.8),
        width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    )
    
    return slide

def create_pie_chart_with_table_slide(prs, title, pie_data, table_headers, table_data):
    """파이 차트 + 표 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, Inches(6), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_format(p, font_size=Pt(28), bold=True, color=COLOR_NAVY)
    
    # matplotlib으로 도넛 차트 생성
    fig, ax = plt.subplots(figsize=(6, 6))
    
    labels = [item['label'] for item in pie_data]
    sizes = [item['value'] for item in pie_data]
    colors = ['#1B3A5C', '#2E5A8C', '#4A90E2', '#7BB3E8', '#AED6F1', '#D4E6F1']
    
    wedges, texts, autotexts = ax.pie(sizes, labels=labels, autopct='%1.1f%%',
                                       colors=colors, startangle=90,
                                       wedgeprops=dict(width=0.5, edgecolor='white'))
    
    # 중앙에 흰색 원 (도넛 스타일)
    centre_circle = plt.Circle((0, 0), 0.3, fc='white')
    ax.add_artist(centre_circle)
    
    plt.tight_layout()
    
    # 차트를 이미지로 저장
    img_buffer = io.BytesIO()
    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close()
    
    # 슬라이드에 추가 (왼쪽 60%)
    chart_width = (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT) * 0.6
    slide.shapes.add_picture(
        img_buffer, MARGIN_LEFT, Inches(1.8),
        width=chart_width
    )
    
    # 표 (오른쪽 40%)
    table_x = MARGIN_LEFT + chart_width + Inches(0.3)
    table_width = (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT) * 0.35
    table_rows = len(table_data) + 1
    table_cols = len(table_headers)
    table_height = Inches(0.4) * table_rows + Inches(0.1)
    table_y = Inches(1.8)
    
    table_shape = slide.shapes.add_table(
        table_rows, table_cols, table_x, table_y, table_width, table_height
    )
    table = table_shape.table
    
    # 헤더 설정
    for i, header in enumerate(table_headers):
        cell = table.cell(0, i)
        cell.text = header
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_format(p, font_size=FONT_SIZE_TABLE_HEADER, 
                           bold=True, color=COLOR_WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
    
    # 데이터 행 설정
    for row_idx, row_data in enumerate(table_data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            
            if col_idx > 0:
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.CENTER
            
            set_paragraph_format(p, font_size=FONT_SIZE_TABLE_BODY, 
                               bold=False, color=COLOR_TEXT)
            
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    
    return slide

def create_donut_chart_slide(prs, title, chart_data, legend_items):
    """도넛 차트 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, Inches(6), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_format(p, font_size=Pt(28), bold=True, color=COLOR_NAVY)
    
    # matplotlib으로 도넛 차트 생성
    fig, ax = plt.subplots(figsize=(8, 6))
    
    labels = [item['label'] for item in chart_data]
    sizes = [item['value'] for item in chart_data]
    colors = [item['color'] for item in chart_data]
    
    wedges, texts, autotexts = ax.pie(sizes, labels=labels, autopct='%1.0f%%',
                                       colors=colors, startangle=90,
                                       wedgeprops=dict(width=0.5, edgecolor='white'))
    
    # 중앙에 흰색 원
    centre_circle = plt.Circle((0, 0), 0.3, fc='white')
    ax.add_artist(centre_circle)
    
    plt.tight_layout()
    
    # 차트를 이미지로 저장
    img_buffer = io.BytesIO()
    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close()
    
    # 슬라이드에 추가
    slide.shapes.add_picture(
        img_buffer, MARGIN_LEFT, Inches(1.8),
        width=Inches(7)
    )
    
    # 범례 (오른쪽에 세로로)
    legend_x = MARGIN_LEFT + Inches(7.5)
    legend_y = Inches(2)
    
    for idx, item in enumerate(legend_items):
        # 아이콘 (작은 사각형)
        icon = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, legend_x, legend_y + idx * Inches(0.9),
            Inches(0.3), Inches(0.3)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = RGBColor(*hex_to_rgb(item['color']))
        icon.line.fill.background()
        
        # 텍스트
        text_box = slide.shapes.add_textbox(
            legend_x + Inches(0.5), legend_y + idx * Inches(0.9),
            Inches(4), Inches(0.7)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = item['text']
        set_paragraph_format(p, font_size=Pt(14), bold=False, color=COLOR_TEXT)
    
    return slide

def hex_to_rgb(hex_color):
    """HEX 색상을 RGB 튜플로 변환"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_line_chart_slide(prs, title, subtitle, x_labels, y_values, footer_text=None):
    """라인 차트 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, Inches(6), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_format(p, font_size=Pt(28), bold=True, color=COLOR_NAVY)
    
    # 부제
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            MARGIN_LEFT, MARGIN_TOP + Inches(0.7), Inches(8), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        set_paragraph_format(p, font_size=Pt(16), bold=False, color=COLOR_TEXT_LIGHT)
    
    # matplotlib으로 라인 차트 생성
    fig, ax = plt.subplots(figsize=(10, 5))
    
    x_positions = range(len(x_labels))
    
    # 라인 플롯
    ax.plot(x_positions, y_values, marker='o', linewidth=3, 
           color='#1B3A5C', markersize=10, markerfacecolor='white', 
           markeredgewidth=2, markeredgecolor='#1B3A5C')
    
    # 데이터 포인트에 값 표시
    for i, v in enumerate(y_values):
        ax.annotate(f'{v}', (i, v), textcoords="offset points", 
                   xytext=(0, 15), ha='center', fontsize=12, 
                   color='#2C3E50', fontweight='bold')
    
    ax.set_xticks(x_positions)
    ax.set_xticklabels(x_labels, fontsize=12)
    ax.set_ylabel('금액 (억원)', fontsize=12)
    ax.grid(axis='y', linestyle='--', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    
    plt.tight_layout()
    
    # 차트를 이미지로 저장
    img_buffer = io.BytesIO()
    plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close()
    
    # 슬라이드에 추가
    slide.shapes.add_picture(
        img_buffer, MARGIN_LEFT, Inches(1.8),
        width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    )
    
    # 푸터 박스
    if footer_text:
        footer_y = Inches(6)
        footer = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            MARGIN_LEFT, footer_y, 
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.8)
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = COLOR_LIGHT_BLUE
        footer.line.fill.background()
        
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_format(p, font_size=Pt(16), bold=True, color=COLOR_NAVY)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return slide

def create_tech_cards_slide(prs, title, subtitle, cards):
    """기술 개발 현황 카드 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 제목
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, Inches(6), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_paragraph_format(p, font_size=Pt(28), bold=True, color=COLOR_NAVY)
    
    # 부제
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            MARGIN_LEFT, MARGIN_TOP + Inches(0.7), Inches(8), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        set_paragraph_format(p, font_size=Pt(16), bold=False, color=COLOR_TEXT_LIGHT)
    
    # 카드 레이아웃
    card_width = (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.4)) / 3
    card_height = Inches(4.5)
    card_start_y = Inches(1.8)
    
    for idx, card in enumerate(cards):
        x = MARGIN_LEFT + idx * (card_width + Inches(0.2))
        
        # 카드 배경
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, card_start_y, card_width, card_height
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = COLOR_WHITE
        card_shape.line.color.rgb = RGBColor(200, 200, 200)
        card_shape.line.width = Pt(1)
        
        # 아이콘 (원형 + 간단한 도형)
        icon_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + card_width/2 - Inches(0.4), 
            card_start_y + Inches(0.3), Inches(0.8), Inches(0.8)
        )
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = COLOR_NAVY
        icon_circle.line.fill.background()
        
        # 아이콘 심볼 (간단한 도형으로 표현)
        if idx == 0:  # AI
            brain = slide.shapes.add_shape(
                MSO_SHAPE.CLOUD, x + card_width/2 - Inches(0.25),
                card_start_y + Inches(0.45), Inches(0.5), Inches(0.5)
            )
            brain.fill.solid()
            brain.fill.fore_color.rgb = COLOR_WHITE
            brain.line.fill.background()
        elif idx == 1:  # 태양광
            sun = slide.shapes.add_shape(
                MSO_SHAPE.SUN, x + card_width/2 - Inches(0.25),
                card_start_y + Inches(0.45), Inches(0.5), Inches(0.5)
            )
            sun.fill.solid()
            sun.fill.fore_color.rgb = COLOR_GOLD
            sun.line.fill.background()
        else:  # 재활용
            recycle = slide.shapes.add_shape(
                MSO_SHAPE.CIRCULAR_ARROW, x + card_width/2 - Inches(0.25),
                card_start_y + Inches(0.45), Inches(0.5), Inches(0.5)
            )
            recycle.fill.solid()
            recycle.fill.fore_color.rgb = COLOR_BLUE
            recycle.line.fill.background()
        
        # 제목
        title_box = slide.shapes.add_textbox(
            x + Inches(0.3), card_start_y + Inches(1.3), 
            card_width - Inches(0.6), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = card['title']
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_format(p, font_size=Pt(18), bold=True, color=COLOR_NAVY)
        
        # 설명
        desc_box = slide.shapes.add_textbox(
            x + Inches(0.3), card_start_y + Inches(2.2), 
            card_width - Inches(0.6), Inches(1.5)
        )
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = card['description']
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_format(p, font_size=Pt(14), bold=False, color=COLOR_TEXT_LIGHT)
        
        # 핵심 수치 (골드 강조)
        if 'highlight' in card:
            hl_box = slide.shapes.add_textbox(
                x + Inches(0.3), card_start_y + Inches(3.8), 
                card_width - Inches(0.6), Inches(0.5)
            )
            tf = hl_box.text_frame
            p = tf.paragraphs[0]
            p.text = card['highlight']
            p.alignment = PP_ALIGN.CENTER
            set_paragraph_format(p, font_size=Pt(16), bold=True, color=COLOR_GOLD)
    
    return slide

def create_conclusion_slide(prs, title, conclusions):
    """결론 슬라이드 생성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경: 네이비 블루
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_NAVY
    bg.line.fill.background()
    
    # 제목 (중앙)
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_format(p, font_size=Pt(32), bold=True, color=COLOR_WHITE)
    
    # 결론 항목들
    content_y = Inches(2.5)
    
    for conclusion in conclusions:
        # 골드 원형 불릿
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, MARGIN_LEFT, content_y + Inches(0.1),
            Inches(0.3), Inches(0.3)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = COLOR_GOLD
        bullet.line.fill.background()
        
        # 텍스트
        text_box = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(0.5), content_y, 
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.5), Inches(0.8)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = conclusion
        set_paragraph_format(p, font_size=Pt(18), bold=False, color=COLOR_WHITE)
        
        content_y += Inches(1.2)
    
    # "감사합니다" (하단 중앙)
    thanks_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(6.2), SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.8)
    )
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "감사합니다"
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_format(p, font_size=Pt(24), bold=True, color=COLOR_GOLD)
    
    return slide


# ============================================================================
# 3. 메인 실행부
# ============================================================================

def main():
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 빈 슬라이드 레이아웃 사용
    blank_layout = prs.slide_layouts[6]
    
    print("슬라이드 생성 중...")
    
    # 1장: 표지
    print("1장: 표지 생성 중...")
    create_title_slide(
        prs, 
        "전국 하수처리장 관리대행 업체 현황 분석",
        "2024 년 시장 동향 및 전망",
        "2026 년 07 월"
    )
    
    # 2장: 목차
    print("2 장: 목차 생성 중...")
    toc_items = [
        "시장 개요", "주요 업체 현황", "매출액 비교",
        "지역별 분포", "처리장 규모별 현황", "입찰 동향",
        "기술 개발 현황", "시장 전망", "결론 및 제언"
    ]
    create_toc_slide(prs, toc_items)
    
    # 3 장: 시장 개요
    print("3 장: 시장 개요 생성 중...")
    market_cards = [
        {'number': '2,458 개소', 'description': '전국 하수처리장 총량\n(2024 년 기준)'},
        {'number': '67.3%', 'description': '민간위탁 비율'},
        {'number': '4.2%', 'description': '연평균 성장률\n(CAGR, 2020-2024)'}
    ]
    create_infographic_cards_slide(
        prs, "시장 개요", market_cards,
        "시장 규모: 약 2 조 8,500 억원 (2024 년 예상)"
    )
    
    # 4 장: 주요 업체 현황
    print("4 장: 주요 업체 현황 생성 중...")
    table_headers = ['업체명', '운영 처리장 수', '처리용량 (톤/일)', '시장점유율 (%)']
    table_data = [
        ['A 환경그룹', '156', '2,850,000', '18.5'],
        ['B 워터테크', '128', '2,100,000', '14.2'],
        ['C 그린시스템', '98', '1,780,000', '11.8'],
        ['D 에코서비스', '87', '1,650,000', '10.5'],
        ['E 클린워터', '72', '1,320,000', '8.9'],
        ['F 물관리공사', '65', '1,180,000', '7.7'],
        ['G 환경산업', '58', '980,000', '6.5'],
        ['H 워터솔루션', '45', '820,000', '5.4']
    ]
    create_table_slide(
        prs, "주요 업체 현황", "상위 8 개사 운영 현황",
        table_headers, table_data,
        "상위 8 개사 합계 시장점유율: 83.5%"
    )
    
    # 5 장: 매출액 현황
    print("5 장: 매출액 현황 생성 중...")
    categories = ['A 환경그룹', 'B 워터테크', 'C 그린시스템', 'D 에코서비스', 'E 클린워터']
    series_2022 = [4200, 3100, 2450, 2180, 1720]
    series_2023 = [4580, 3420, 2680, 2350, 1890]
    series_2024 = [4950, 3680, 2890, 2540, 2050]
    create_bar_chart_slide(
        prs, "매출액 현황", "상위 5 개사 최근 3 개년 매출액 비교",
        categories, 
        [series_2022, series_2023, series_2024],
        ['2022 년', '2023 년', '2024 년']
    )
    
    # 6 장: 지역별 업체 분포
    print("6 장: 지역별 업체 분포 생성 중...")
    pie_data = [
        {'label': '수도권', 'value': 32.5},
        {'label': '영남권', 'value': 20.9},
        {'label': '호남권', 'value': 16.4},
        {'label': '충청권', 'value': 13.3},
        {'label': '강원권', 'value': 5.9},
        {'label': '기타', 'value': 11.0}
    ]
    region_table_headers = ['권역', '업체 수', '비중 (%)']
    region_table_data = [
        ['수도권', '485', '32.5'],
        ['영남권', '312', '20.9'],
        ['호남권', '245', '16.4'],
        ['충청권', '198', '13.3'],
        ['강원권', '89', '5.9'],
        ['기타', '169', '11.0']
    ]
    create_pie_chart_with_table_slide(
        prs, "지역별 업체 분포",
        pie_data, region_table_headers, region_table_data
    )
    
    # 7 장: 처리장 규모별 현황
    print("7 장: 처리장 규모별 현황 생성 중...")
    donut_data = [
        {'label': '대규모\n(10,000 톤/일 이상)', 'value': 15, 'color': '#1B3A5C'},
        {'label': '중규모\n(1,000-10,000 톤/일)', 'value': 45, 'color': '#4A90E2'},
        {'label': '소규모\n(1,000 톤/일 미만)', 'value': 40, 'color': '#AED6F1'}
    ]
    donut_legend = [
        {'color': '#1B3A5C', 'text': '대규모 (10,000 톤/일 이상) - 15%'},
        {'color': '#4A90E2', 'text': '중규모 (1,000-10,000 톤/일) - 45%'},
        {'color': '#AED6F1', 'text': '소규모 (1,000 톤/일 미만) - 40%'}
    ]
    create_donut_chart_slide(prs, "처리장 규모별 현황", donut_data, donut_legend)
    
    # 8 장: 입찰 동향
    print("8 장: 입찰 동향 생성 중...")
    years = ['2020', '2021', '2022', '2023', '2024']
    bid_amounts = [420, 485, 520, 580, 625]
    create_line_chart_slide(
        prs, "입찰 동향", "최근 5 개년 평균 입찰 금액 추이",
        years, bid_amounts,
        "평균 낙찰률: 92.3%"
    )
    
    # 9 장: 기술 개발 현황
    print("9 장: 기술 개발 현황 생성 중...")
    tech_cards = [
        {
            'title': 'AI 기반 운영 최적화',
            'description': '에너지 절감 25%, 운영 효율 30% 향상',
            'highlight': ''
        },
        {
            'title': '에너지 자립화',
            'description': '태양광+바이오가스, 자립률 45% 달성',
            'highlight': ''
        },
        {
            'title': '슬러지 감량화',
            'description': '부피 60% 감소, 재활용률 80%',
            'highlight': ''
        }
    ]
    create_tech_cards_slide(prs, "기술 개발 현황", "주요 기술 트렌드 3 가지", tech_cards)
    
    # 10 장: 결론 및 제언
    print("10 장: 결론 및 제언 생성 중...")
    conclusions = [
        "시장 안정성: 연 4.2% 성장으로 안정적 투자처",
        "기술 혁신: AI·에너지 자립화 기술이 경쟁력 결정",
        "투자 유망 분야: 스마트 워터관리, 에너지 효율화, 슬러지 자원화"
    ]
    create_conclusion_slide(prs, "결론 및 제언", conclusions)
    
    # 저장
    output_file = '하수처리장_관리대행_업체_현황_v3.pptx'
    prs.save(output_file)
    print(f"\n✅ 프레젠테이션이 성공적으로 생성되었습니다: {output_file}")
    print(f"📊 총 {len(prs.slides)}개의 슬라이드가 포함되었습니다.")

if __name__ == '__main__':
    main()
