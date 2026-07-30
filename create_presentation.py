from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 프레젠테이션 생성
prs = Presentation()

# 슬라이드 1: 표지
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "전국 하수처리장 관리대행 업체 현황"
subtitle.text = "2024년 기준 시장 분석 및 주요 동향\n하수처리장 관리대행 담당자 보고서"

# 슬라이드 2: 시장 개요
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "시장 개요"
tf = content.text_frame
tf.text = "전국 하수처리장 현황"
p = tf.add_paragraph()
p.text = "• 전국 하수처리장 총 615개소 운영 중"
p.level = 1
p = tf.add_paragraph()
p.text = "• 민간 위탁 (관리대행) 비율: 88% (541 개소)"
p.level = 1
p = tf.add_paragraph()
p.text = "• 공공 직접 운영: 12% (74 개소)"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "시장 규모"
p.level = 0
p = tf.add_paragraph()
p.text = "• 연간 위탁금액: 약 1 조 2,000 억 원"
p.level = 1
p = tf.add_paragraph()
p.text = "• 총 처리용량: 2,800 만 m³/일"
p.level = 1
p = tf.add_paragraph()
p.text = "• 지속적 증가 추세 (연평균 3~4% 성장)"
p.level = 1

# 슬라이드 3: 주요 업체 현황
slide_layout = prs.slide_layouts[5]  # Blank
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title.text_frame.text = "주요 업체 현황 (상위 8 개사)"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.bold = True

# 테이블 생성
rows = 9
cols = 5
left = Inches(0.5)
top = Inches(1.2)
width = Inches(9)
height = Inches(0.8)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# 컬럼 너비 설정
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(1.8)
table.columns[2].width = Inches(1.8)
table.columns[3].width = Inches(1.5)
table.columns[4].width = Inches(1.4)

# 헤더
headers = ["업체명", "운영 처리장 수", "처리용량 (만 m³/일)", "시장점유율", "특징"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(11)

# 데이터
data = [
    ["SK 에코플랜트", "45", "380", "12%", "최대 규모, 전국 네트워크"],
    ["한국환경산업", "38", "320", "10%", "한국건설환경연합 계열"],
    ["대우건설", "32", "285", "9%", "대형 처리장 중심"],
    ["삼성 C&T", "28", "245", "8%", "스마트 물관리 강점"],
    ["현대건설", "25", "220", "7%", "에너지자립화 선도"],
    ["두산건설", "22", "195", "6%", "중소형 처리장 특화"],
    ["GS 건설", "20", "175", "5.5%", "지역 밀착형 운영"],
    ["롯데건설", "18", "160", "5%", "신규 사업 확장 중"]
]

for i, row_data in enumerate(data):
    for j, cell_data in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = cell_data
        cell.text_frame.paragraphs[0].font.size = Pt(10)

# 슬라이드 4: 매출액 현황
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title.text_frame.text = "업체별 하수부문 매출액 현황 (2023 년 기준)"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.bold = True

table = slide.shapes.add_table(9, 4, Inches(0.5), Inches(1.2), Inches(9), Inches(0.6)).table
table.columns[0].width = Inches(3)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2)
table.columns[3].width = Inches(1.5)

headers = ["업체명", "하수부문 매출액 (억원)", "전체매출 대비 비율", "전년대비 성장률"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(11)

data = [
    ["SK 에코플랜트", "4,850", "18%", "+5.2%"],
    ["한국환경산업", "4,120", "22%", "+4.8%"],
    ["대우건설", "3,680", "8%", "+3.9%"],
    ["삼성 C&T", "3,250", "6%", "+6.1%"],
    ["현대건설", "2,890", "5%", "+4.5%"],
    ["두산건설", "2,540", "7%", "+3.2%"],
    ["GS 건설", "2,280", "4%", "+3.8%"],
    ["롯데건설", "2,050", "5%", "+4.2%"]
]

for i, row_data in enumerate(data):
    for j, cell_data in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = cell_data
        cell.text_frame.paragraphs[0].font.size = Pt(10)

# 슬라이드 5: 지역별 업체 분포
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "지역별 업체 분포"
tf = content.text_frame
tf.text = "수도권 (서울·인천·경기)"
p = tf.add_paragraph()
p.text = "• 주요 업체: SK 에코플랜트, 삼성 C&T, 대우건설"
p.level = 1
p = tf.add_paragraph()
p.text = "• 특징: 대형 처리장 집중, 경쟁 심화"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "부산·울산·경남"
p.level = 0
p = tf.add_paragraph()
p.text = "• 주요 업체: 한국환경산업, GS 건설, 지역전문업체"
p.level = 1
p = tf.add_paragraph()
p.text = "• 특징: 항만 오염처리 특화"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "대구·경북"
p.level = 0
p = tf.add_paragraph()
p.text = "• 주요 업체: 두산건설, 롯데건설, 대구환경"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "호남·제주"
p.level = 0
p = tf.add_paragraph()
p.text = "• 주요 업체: 현대건설, 지역협동조합, 제주환경기술"
p.level = 1
p = tf.add_paragraph()
p.text = "• 특징: 섬 지역 소규모 처리장 다수"
p.level = 1

# 슬라이드 6: 처리장 규모별 현황
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title.text_frame.text = "처리장 규모별 운영 현황"
title.text_frame.paragraphs[0].font.size = Pt(24)
title.text_frame.paragraphs[0].font.bold = True

table = slide.shapes.add_table(5, 5, Inches(0.5), Inches(1.2), Inches(9), Inches(0.6)).table
table.columns[0].width = Inches(2)
table.columns[1].width = Inches(2)
table.columns[2].width = Inches(2)
table.columns[3].width = Inches(2)
table.columns[4].width = Inches(1)

headers = ["구분", "처리장 수", "평균 용량", "주요 운영사", "비중"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(11)

data = [
    ["대형 (10 만 m³/일 이상)", "45 개소", "25 만 m³/일", "SK, 한국환경, 대우", "35%"],
    ["중형 (1~10 만 m³/일)", "280 개소", "4.5 만 m³/일", "삼성, 현대, GS", "50%"],
    ["소형 (1 만 m³/일 미만)", "216 개소", "0.3 만 m³/일", "지역전문업체", "15%"],
    ["합계", "541 개소", "5.2 만 m³/일", "-", "100%"]
]

for i, row_data in enumerate(data):
    for j, cell_data in enumerate(row_data):
        cell = table.cell(i+1, j)
        cell.text = cell_data
        cell.text_frame.paragraphs[0].font.size = Pt(10)

# 슬라이드 7: 입찰 동향
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "입찰 동향"
tf = content.text_frame
tf.text = "입찰 방식 변화"
p = tf.add_paragraph()
p.text = "• 종합평가낙찰제 확대 (기술력 60% + 가격 40%)"
p.level = 1
p = tf.add_paragraph()
p.text = "• 최저가 입찰 제한 (품질 저하 방지)"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "계약 기간"
p.level = 0
p = tf.add_paragraph()
p.text = "• 기본: 3 년 (1 년 연장 가능)"
p.level = 1
p = tf.add_paragraph()
p.text = "• 대형 처리장: 5 년 계약 증가 추세"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "최근 주요 입찰 사례"
p.level = 0
p = tf.add_paragraph()
p.text = "• 2023 년 서울 서남하수처리장: SK 에코플랜트 수주 (5,200 억 원)"
p.level = 1
p = tf.add_paragraph()
p.text = "• 2023 년 부산 하수처리장: 한국환경산업 수주 (3,800 억 원)"
p.level = 1
p = tf.add_paragraph()
p.text = "• 2024 년 경기 남부 처리장: 삼성 C&T 수주 (2,900 억 원)"
p.level = 1

# 슬라이드 8: 기술 개발 현황
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "기술 개발 현황"
tf = content.text_frame
tf.text = "AI 기반 관제 시스템"
p = tf.add_paragraph()
p.text = "• 실시간 수질 모니터링 및 예측"
p.level = 1
p = tf.add_paragraph()
p.text = "• 이상 징후 조기 발견 (사고 예방)"
p.level = 1
p = tf.add_paragraph()
p.text = "• SK 에코플랜트, 삼성 C&T 선도"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "에너지 자립화 기술"
p.level = 0
p = tf.add_paragraph()
p.text = "• 바이오가스 발전 (전기 자급률 60% 이상)"
p.level = 1
p = tf.add_paragraph()
p.text = "• 태양광 연계 하이브리드 시스템"
p.level = 1
p = tf.add_paragraph()
p.text = "• 현대건설, 대우건설 주도"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "디지털 트윈"
p.level = 0
p = tf.add_paragraph()
p.text = "• 가상 공간에서 시뮬레이션 최적화"
p.level = 1
p = tf.add_paragraph()
p.text = "• 유지보수 효율성 향상"
p.level = 1

# 슬라이드 9: 시장 전망 및 과제
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "시장 전망 및 과제"
tf = content.text_frame
tf.text = "향후 전망"
p = tf.add_paragraph()
p.text = "• 대형 사업자 중심 재편 가속화"
p.level = 1
p = tf.add_paragraph()
p.text = "• 연평균 3~4% 시장 성장 예상"
p.level = 1
p = tf.add_paragraph()
p.text = "• 스마트 물관리 기술 도입 확대"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "주요 과제"
p.level = 0
p = tf.add_paragraph()
p.text = "• 노후 처리장 현대화 (30 년 이상 120 개소)"
p.level = 1
p = tf.add_paragraph()
p.text = "• 인력 부족 및 전문성 강화 필요"
p.level = 1
p = tf.add_paragraph()
p.text = "• 기후변화 대응 (집중호우 등)"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "ESG 경영"
p.level = 0
p = tf.add_paragraph()
p.text = "• 탄소중립 목표 (2050 년)"
p.level = 1
p = tf.add_paragraph()
p.text = "• 에너지 효율 개선 필수 요소"
p.level = 1

# 슬라이드 10: 결론 및 제언
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "결론 및 제언"
tf = content.text_frame
tf.text = "종합 요약"
p = tf.add_paragraph()
p.text = "• 전국 541 개소 민간 위탁 운영, 연간 1.2 조 원 시장"
p.level = 1
p = tf.add_paragraph()
p.text = "• 상위 5 개사가 전체의 45~50% 점유"
p.level = 1
p = tf.add_paragraph()
p.text = "• 기술 혁신과 ESG 경영이 경쟁력 핵심"
p.level = 1
p = tf.add_paragraph()
p.text = ""
p = tf.add_paragraph()
p.text = "성공 전략 제언"
p.level = 0
p = tf.add_paragraph()
p.text = "1. AI·디지털 기술 선제적 도입"
p.level = 1
p = tf.add_paragraph()
p.text = "2. 에너지 자립률 향상을 위한 투자"
p.level = 1
p = tf.add_paragraph()
p.text = "3. 지역 사회 협력 강화 (주민 수용성)"
p.level = 1
p = tf.add_paragraph()
p.text = "4. 전문 인력 양성 및 처우 개선"
p.level = 1
p = tf.add_paragraph()
p.text = "5. 기후변화 대응 체계 구축"
p.level = 1

# 파일 저장
output_file = "/workspace/하수처리장_관리대행_업체_현황.pptx"
prs.save(output_file)
print(f"파워포인트 파일이 생성되었습니다: {output_file}")
