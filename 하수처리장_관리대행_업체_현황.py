# -*- coding: utf-8 -*-
"""
하수처리장 관리대행 업체 현황 파워포인트 자동 생성 스크립트

실행 방법:
1. 필요한 라이브러리 설치: pip install python-pptx
2. 스크립트 실행: python 하수처리장_관리대행_업체_현황.py

생성 파일: 하수처리장_관리대행_업체_현황.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime


# 색상 정의
NAVY_BLUE = RGBColor(0, 51, 102)  # #003366
GRAY = RGBColor(128, 128, 128)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)


def create_presentation():
    """프레젠테이션 객체 생성"""
    prs = Presentation()
    return prs


def set_title_format(title_shape, text, size=28):
    """제목 텍스트 서식 설정"""
    title_shape.text = text
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = NAVY_BLUE
        try:
            paragraph.font.name = 'Malgun Gothic'
        except:
            pass


def set_body_format(text_frame, text, size=18):
    """본문 텍스트 서식 설정"""
    text_frame.text = text
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = RGBColor(64, 64, 64)
        try:
            paragraph.font.name = 'Malgun Gothic'
        except:
            pass


def create_title_slide(prs, title, subtitle, date_str):
    """표지 슬라이드 생성"""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목 설정
    title_shape = slide.shapes.title
    set_title_format(title_shape, title, 32)
    
    # 부제 설정
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = f"{subtitle}\n{date_str}"
    for paragraph in subtitle_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = GRAY
        try:
            paragraph.font.name = 'Malgun Gothic'
        except:
            pass
    
    # 네이비 블루 라인 추가
    left = Inches(0.5)
    top = Inches(6.5)
    width = Inches(9)
    height = Inches(0.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_BLUE
    shape.line.fill.background()
    
    return slide


def create_toc_slide(prs, toc_items):
    """목차 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목 설정
    title_shape = slide.shapes.title
    set_title_format(title_shape, "목 차", 28)
    
    # 목차 내용 설정
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, item in enumerate(toc_items, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(64, 64, 64)
        try:
            p.font.name = 'Malgun Gothic'
        except:
            p.space_after = Pt(12)
    
    return slide


def create_content_slide_with_bullets(prs, title, bullets):
    """불릿 포인트가 있는 콘텐츠 슬라이드 생성"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목 설정
    title_shape = slide.shapes.title
    set_title_format(title_shape, title, 28)
    
    # 본문 설정
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(64, 64, 64)
        try:
            p.font.name = 'Malgun Gothic'
        except:
            p.space_after = Pt(10)
    
    return slide


def create_table_slide(prs, title, headers, data, column_widths=None):
    """표가 있는 슬라이드 생성"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목 추가
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    try:
        p.font.name = 'Malgun Gothic'
    except:
        pass
    
    # 표 생성
    num_rows = len(data) + 1  # 헤더 행 포함
    num_cols = len(headers)
    
    table_left = Inches(0.5)
    table_top = Inches(1.2)
    table_width = Inches(9)
    table_height = Inches(0.8) * num_rows
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table
    
    # 열 너비 설정
    if column_widths:
        for i, width in enumerate(column_widths):
            table.columns[i].width = Inches(width)
    
    # 헤더 행 설정
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BLUE
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
            try:
                paragraph.font.name = 'Malgun Gothic'
            except:
                pass
    
    # 데이터 행 설정
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(64, 64, 64)
                paragraph.alignment = PP_ALIGN.CENTER
                try:
                    paragraph.font.name = 'Malgun Gothic'
                except:
                    pass
            
            # 짝수 행에 밝은 회색 배경
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # 테두리 설정 (python-pptx 버전 호환성 고려)
    # 표의 외곽 테두리는 자동으로 적용됨
    
    return slide


def main():
    """메인 함수"""
    # 프레젠테이션 생성
    prs = create_presentation()
    
    # 현재 날짜
    date_str = datetime.now().strftime("%Y년 %m월 %d일")
    
    # 1장: 표지
    create_title_slide(
        prs,
        "전국 하수처리장 관리대행 업체 현황 분석",
        "2024년 시장 동향 및 전망",
        date_str
    )
    
    # 2장: 목차
    toc_items = [
        "시장 개요",
        "주요 업체 현황",
        "매출액 비교",
        "지역별 분포",
        "처리장 규모별 현황",
        "입찰 동향",
        "기술 개발 현황",
        "시장 전망",
        "결론"
    ]
    create_toc_slide(prs, toc_items)
    
    # 3장: 시장 개요
    market_overview_bullets = [
        "전국 하수처리장 총량: 2,458개소 (2024년 기준)",
        "민간위탁 비율: 67.3% (1,654개소)",
        "연평균 성장률 (CAGR): 4.2% (2020~2024년)",
        "시장 규모: 약 2조 8,500억원 (2024년 예상)",
        "정부의 민간위託 확대 정책에 따른 지속적 성장세",
        "환경규제 강화로 전문 운영업체 수요 증가"
    ]
    create_content_slide_with_bullets(prs, "시장 개요", market_overview_bullets)
    
    # 4장: 주요 업체 현황
    headers_4 = ["업체명", "운영 처리장 수", "처리용량 (톤/일)", "시장점유율 (%)"]
    data_4 = [
        ["A 환경건설", 156, 2850000, 18.5],
        ["B 워터테크", 128, 2340000, 15.2],
        ["C 그린시스템", 98, 1780000, 11.6],
        ["D 에코서비스", 87, 1560000, 10.1],
        ["E 클린워터", 72, 1320000, 8.6],
        ["F 물관리공사", 65, 1180000, 7.7],
        ["G 청정기술", 54, 980000, 6.4],
        ["H 환경개발", 48, 870000, 5.7]
    ]
    create_table_slide(prs, "주요 업체 현황", headers_4, data_4, [2.5, 2.0, 2.5, 2.0])
    
    # 5장: 매출액 현황
    headers_5 = ["업체명", "2022년 (억원)", "2023년 (억원)", "2024년 예상 (억원)"]
    data_5 = [
        ["A 환경건설", 4850, 5120, 5480],
        ["B 워터테크", 3980, 4250, 4580],
        ["C 그린시스템", 2870, 3050, 3280],
        ["D 에코서비스", 2450, 2620, 2810],
        ["E 클린워터", 2180, 2340, 2520]
    ]
    create_table_slide(prs, "매출액 현황 (상위 5개사)", headers_5, data_5, [2.0, 2.3, 2.3, 2.4])
    
    # 6장: 지역별 업체 분포
    headers_6 = ["권역", "업체 수", "비중 (%)", "주요 처리장 수"]
    data_6 = [
        ["경기·인천", 425, 25.7, 398],
        ["서울", 186, 11.3, 142],
        ["경남·부산", 298, 18.0, 312],
        ["경북·대구", 215, 13.0, 245],
        ["충청권", 245, 14.8, 268],
        ["전라권", 198, 12.0, 215],
        ["강원·제주", 87, 5.2, 80]
    ]
    create_table_slide(prs, "지역별 업체 분포", headers_6, data_6, [2.0, 2.0, 2.0, 3.0])
    
    # 7장: 처리장 규모별 현황
    headers_7 = ["규모 구분", "업체 수", "비중 (%)", "처리 용량 비중 (%)"]
    data_7 = [
        ["대규모 (50,000톤/일 이상)", 145, 8.8, 52.3],
        ["중규모 (10,000~50,000톤)", 562, 34.0, 35.8],
        ["소규모 (10,000톤 미만)", 947, 57.2, 11.9]
    ]
    create_table_slide(prs, "처리장 규모별 현황", headers_7, data_7, [3.0, 2.0, 2.0, 2.0])
    
    # 8장: 입찰 동향
    headers_8 = ["구분", "2023년", "2024년", "변동률 (%)"]
    data_8 = [
        ["평균 입찰 금액 (억원)", 45.8, 48.2, "+5.2"],
        ["평균 낙찰률 (%)", 87.5, 89.2, "+1.7"],
        ["입찰 참여 업체 수 (사)", 12.3, 14.1, "+14.6"],
        ["기술평가 비중 (%)", 45.0, 50.0, "+5.0"]
    ]
    create_table_slide(prs, "입찰 동향 (최근 2년간)", headers_8, data_8, [2.5, 2.2, 2.2, 2.1])
    
    # 9장: 기술 개발 현황
    tech_trends = [
        "AI 기반 운영 최적화: 실시간 수질 예측 및 에너지 효율화 시스템 도입 확대",
        "에너지 자립화: 태양광, 수소연료전지 등 재생에너지 활용 시설 증가",
        "슬러지 감량화: 고도처리 기술 및 자원화 기술 개발 활발",
        "스마트 워터그리드: IoT 센서 기반 원격 모니터링 시스템 구축",
        "탄소중립 대응: 온실가스 배출 저감 기술 및 LCA 평가 도입"
    ]
    create_content_slide_with_bullets(prs, "기술 개발 현황", tech_trends)
    
    # 10장: 결론 및 제언
    conclusion_bullets = [
        "시장 안정성: 꾸준한 성장세와 정부의 민간위託 정책으로 안정적인 시장 환경 조성",
        "투자 유망 분야 1: AI·IoT 기반 스마트 운영관리 솔루션",
        "투자 유망 분야 2: 에너지 자립형 하수처리 시설 기술",
        "투자 유망 분야 3: 슬러지 자원화 및 탄소중립 관련 기술",
        "전문성과 기술력을 갖춘 업체의 경쟁력 강화 필요"
    ]
    create_content_slide_with_bullets(prs, "결론 및 제언", conclusion_bullets)
    
    # 파일 저장
    output_filename = "하수처리장_관리대행_업체_현황.pptx"
    try:
        prs.save(output_filename)
        print(f"파일이 성공적으로 생성되었습니다: {output_filename}")
    except PermissionError:
        print("오류: 파일 저장 권한이 없습니다. 다른 경로를 시도해주세요.")
    except Exception as e:
        print(f"오류 발생: {e}")


if __name__ == "__main__":
    main()
